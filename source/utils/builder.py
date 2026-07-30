"""
Appendix builder - assembles a PTG-formatted PPTX from a project config,
sheet list, and a folder of exported images.
"""

import os

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.util import Mm, Pt

from utils.config_loader import find_image
from utils.titleblock_constants import (
    CONTENT_H,
    CONTENT_L,
    CONTENT_T,
    CONTENT_W,
    SLIDE_H_MM,
    SLIDE_W_MM,
    WHITE,
)
from utils.titleblock_renderer import draw_titleblock


def _remove_all_slides(prs):
    """Remove all slides from a presentation while preserving theme/master data."""
    slide_ids = list(prs.slides._sldIdLst)
    for slide_id in slide_ids:
        rel_id = slide_id.rId
        prs.part.drop_rel(rel_id)
        prs.slides._sldIdLst.remove(slide_id)


def _get_blank_layout(prs):
    """Return a blank slide layout if available, otherwise fall back to the last layout."""
    for layout in prs.slide_layouts:
        if layout.name.lower() == "blank":
            return layout
    return prs.slide_layouts[len(prs.slide_layouts) - 1]


def _fit_image_in_box(img_w_px, img_h_px, box_w_mm, box_h_mm):
    """
    Calculate image placement (left_mm, top_mm, w_mm, h_mm) to fit the image
    inside the box while preserving aspect ratio and centering it.
    """
    img_ratio = img_w_px / img_h_px
    box_ratio = box_w_mm / box_h_mm

    if img_ratio > box_ratio:
        w = box_w_mm
        h = box_w_mm / img_ratio
    else:
        h = box_h_mm
        w = box_h_mm * img_ratio

    left = CONTENT_L + (box_w_mm - w) / 2
    top = CONTENT_T + (box_h_mm - h) / 2
    return left, top, w, h


def _get_image_size(image_path):
    """Return (width_px, height_px) for an image file."""
    from PIL import Image

    with Image.open(image_path) as img:
        return img.size


def build_appendix(
    exports_dir: str,
    output_path: str,
    project_config: dict,
    sheets: list,
    logo_path: str = None,
    template_path: str = None,
    on_progress=None,
) -> dict:
    """
    Build the appendix PPTX and return a summary dict.
    """
    prs = (
        Presentation(template_path)
        if template_path and os.path.isfile(template_path)
        else Presentation()
    )
    prs.slide_width = Mm(SLIDE_W_MM)
    prs.slide_height = Mm(SLIDE_H_MM)

    _remove_all_slides(prs)
    blank_layout = _get_blank_layout(prs)

    missing = []
    slides_built = 0
    total = len(sheets)

    for i, sheet in enumerate(sheets):
        if on_progress:
            message = f"Building sheet {sheet.get('sheet_number', '?')} - {sheet.get('drawing_title_1', '')}"
            on_progress(i + 1, total, message)

        slide = prs.slides.add_slide(blank_layout)

        bg = slide.background
        fill = bg.fill
        fill.solid()
        fill.fore_color.rgb = WHITE

        fields = dict(project_config)
        fields.update({k: v for k, v in sheet.items() if k != "filename_pattern"})

        draw_titleblock(slide, fields, logo_path=logo_path)

        pattern = sheet.get("filename_pattern", "")
        source_path = sheet.get("source_path", "")
        img_path = source_path if source_path and os.path.isfile(source_path) else None
        if not img_path and pattern:
            img_path = find_image(exports_dir, pattern)

        if img_path and os.path.isfile(img_path):
            try:
                w_px, h_px = _get_image_size(img_path)
                left, top, w, h = _fit_image_in_box(w_px, h_px, CONTENT_W, CONTENT_H)
                slide.shapes.add_picture(img_path, Mm(left), Mm(top), Mm(w), Mm(h))
            except Exception as exc:
                _add_missing_placeholder(slide, sheet, f"Unable to read image: {exc}")
                missing.append(
                    (sheet.get("sheet_number"), sheet.get("drawing_title_1"), pattern)
                )
        else:
            missing_reason = (
                f"No file found matching '{pattern}'"
                if pattern
                else "No source image selected"
            )
            _add_missing_placeholder(slide, sheet, missing_reason)
            missing.append(
                (
                    sheet.get("sheet_number"),
                    sheet.get("drawing_title_1"),
                    pattern or source_path,
                )
            )

        slides_built += 1

    if on_progress:
        on_progress(total, total, "Saving...")

    os.makedirs(os.path.dirname(os.path.abspath(output_path)), exist_ok=True)
    prs.save(output_path)

    return {
        "slides_built": slides_built,
        "missing": missing,
        "output_path": output_path,
    }


def _add_missing_placeholder(slide, sheet, reason):
    """Draw a grey placeholder box with a message when an image is missing."""
    from lxml import etree
    from pptx.enum.text import PP_ALIGN
    from pptx.oxml.ns import qn

    shape = slide.shapes.add_shape(
        1, Mm(CONTENT_L), Mm(CONTENT_T), Mm(CONTENT_W), Mm(CONTENT_H)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = RGBColor(0xF0, 0xF0, 0xF0)
    shape.line.color.rgb = RGBColor(0xCC, 0xCC, 0xCC)

    tf = shape.text_frame
    tf.word_wrap = True

    title_para = tf.paragraphs[0]
    title_para.alignment = PP_ALIGN.CENTER
    run = title_para.add_run()
    run.text = f"[ {sheet.get('drawing_title_1', 'Image')} ]"
    run.font.size = Pt(18)
    run.font.bold = True
    run.font.color.rgb = RGBColor(0x99, 0x99, 0x99)
    run.font.name = "Arial"

    paragraph = etree.SubElement(tf._txBody, qn("a:p"))
    run2 = etree.SubElement(paragraph, qn("a:r"))
    run_props = etree.SubElement(run2, qn("a:rPr"), lang="en-AU")
    run_props.set("sz", "1000")
    text = etree.SubElement(run2, qn("a:t"))
    text.text = reason
