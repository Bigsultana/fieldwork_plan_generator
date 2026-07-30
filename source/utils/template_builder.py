"""
Utilities for generating a reusable appendix PowerPoint template.
"""

import os

from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE
from pptx.util import Mm, Pt

from utils.titleblock_constants import (
    CONTENT_H,
    CONTENT_L,
    CONTENT_T,
    CONTENT_W,
    SLIDE_H_MM,
    SLIDE_W_MM,
    WHITE,
)
from utils.titleblock_renderer import draw_titleblock


def _get_blank_layout(prs):
    for layout in prs.slide_layouts:
        if layout.name.lower() == "blank":
            return layout
    return prs.slide_layouts[len(prs.slide_layouts) - 1]


def create_appendix_template(base_pptx_path: str, output_path: str) -> str:
    """
    Build a reusable appendix template PPTX using the supplied base presentation.
    """
    prs = (
        Presentation(base_pptx_path)
        if base_pptx_path and os.path.isfile(base_pptx_path)
        else Presentation()
    )
    prs.slide_width = Mm(SLIDE_W_MM)
    prs.slide_height = Mm(SLIDE_H_MM)

    blank_layout = _get_blank_layout(prs)
    slide = prs.slides.add_slide(blank_layout)

    bg = slide.background
    bg.fill.solid()
    bg.fill.fore_color.rgb = WHITE

    placeholder_fields = {
        "project_title": "PROJECT TITLE",
        "project_address": "PROJECT ADDRESS",
        "project_number": "PTG-YYYY-XXX",
        "client_name": "CLIENT NAME",
        "drawing_title_1": "DRAWING TITLE",
        "drawing_title_2": "SUBTITLE",
        "drawing_title_3": "",
        "drawn_by": "XX",
        "designed_by": "XX",
        "approved_by": "XX",
        "date": "DD.MM.YY",
        "scale": "NTS",
        "sheet_number": "001",
        "sheet_prefix": "A",
        "revision": "-",
        "drawing_status": "FOR INFORMATION",
        "company_name": "PTG CONSULTING",
        "company_address": "Level 3, 159 Coronation Drive, Milton QLD 4064",
        "company_phone": "(07) 3444 6666",
        "company_email": "admin@ptgconsulting.com.au",
        "company_website": "www.ptgconsulting.com.au",
    }

    draw_titleblock(slide, placeholder_fields, logo_path=None)

    # Show the intended image zone so the template is visually self-explanatory.
    frame = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Mm(CONTENT_L),
        Mm(CONTENT_T),
        Mm(CONTENT_W),
        Mm(CONTENT_H),
    )
    frame.fill.background()
    frame.line.dash_style = None
    frame.line.width = Pt(1)

    tf = frame.text_frame
    tf.word_wrap = True
    tf.clear()
    para = tf.paragraphs[0]
    run = para.add_run()
    run.text = "IMAGE CONTENT AREA\n\nBuilder places exported figures here."
    run.font.name = "Arial"
    run.font.size = Pt(22)
    run.font.bold = True

    notes = slide.shapes.add_textbox(Mm(18), Mm(18), Mm(120), Mm(18))
    notes_tf = notes.text_frame
    notes_tf.word_wrap = True
    notes_tf.clear()
    p = notes_tf.paragraphs[0]
    r = p.add_run()
    r.text = "PTG appendix template generated from DXF-derived titleblock geometry."
    r.font.name = "Arial"
    r.font.size = Pt(10)

    os.makedirs(os.path.dirname(os.path.abspath(output_path)), exist_ok=True)
    prs.save(output_path)
    return output_path
