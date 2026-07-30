"""
Titleblock renderer - draws the PTG Consulting A1 titleblock onto a python-pptx slide.
All geometry sourced from DXF parse of A1_HORIZONTAL_-_PTG_CONSULTING_COMMERCIAL.
"""

import os

from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Mm, Pt

from utils import titleblock_constants as tc


def _mm(value):
    return Mm(value)


def _add_line(slide, x1_mm, y1_mm, x2_mm, y2_mm, width_pt=0.5, color=tc.BLACK):
    if abs(y1_mm - y2_mm) < 0.1:
        left = _mm(min(x1_mm, x2_mm))
        top = _mm(y1_mm) - Pt(width_pt) / 2
        width = _mm(abs(x2_mm - x1_mm))
        height = Pt(width_pt)
        shape = slide.shapes.add_shape(1, left, top, width, height)
    elif abs(x1_mm - x2_mm) < 0.1:
        left = _mm(x1_mm) - Pt(width_pt) / 2
        top = _mm(min(y1_mm, y2_mm))
        width = Pt(width_pt)
        height = _mm(abs(y2_mm - y1_mm))
        shape = slide.shapes.add_shape(1, left, top, width, height)
    else:
        return

    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()


def _add_rect(
    slide, left_mm, top_mm, w_mm, h_mm, fill=None, line_color=tc.BLACK, line_pt=0.5
):
    shape = slide.shapes.add_shape(1, _mm(left_mm), _mm(top_mm), _mm(w_mm), _mm(h_mm))
    if fill is not None:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill
    else:
        shape.fill.background()
    shape.line.color.rgb = line_color
    shape.line.width = Pt(line_pt)
    return shape


def _fit_pt(text, base_pt, min_pt, width_mm, height_mm):
    if not text:
        return Pt(base_pt)
    size = float(base_pt)
    width_pt = max(width_mm * 2.83465, 1)
    height_pt = max(height_mm * 2.83465, 1)
    length = max(len(text), 1)
    while size > min_pt:
        chars_per_line = max(int(width_pt / max(size * 0.50, 1)), 1)
        lines = max(1, (length + chars_per_line - 1) // chars_per_line)
        if lines * size * 1.10 <= height_pt:
            break
        size -= 0.5
    return Pt(max(size, min_pt))


def _add_textbox(
    slide,
    left_mm,
    top_mm,
    w_mm,
    h_mm,
    text,
    font_size,
    bold=False,
    italic=False,
    color=tc.BLACK,
    align=PP_ALIGN.LEFT,
    word_wrap=False,
):
    tx_box = slide.shapes.add_textbox(_mm(left_mm), _mm(top_mm), _mm(w_mm), _mm(h_mm))
    tf = tx_box.text_frame
    tf.word_wrap = word_wrap
    tf.margin_left = Mm(0.35)
    tf.margin_right = Mm(0.35)
    tf.margin_top = Mm(0.15)
    tf.margin_bottom = Mm(0.15)
    paragraph = tf.paragraphs[0]
    paragraph.alignment = align
    run = paragraph.add_run()
    run.text = text
    run.font.size = font_size
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    run.font.name = "Arial"
    return tx_box


def _fitted_textbox(
    slide,
    left_mm,
    top_mm,
    w_mm,
    h_mm,
    text,
    base_pt,
    min_pt=6.0,
    bold=False,
    color=tc.BLACK,
    align=PP_ALIGN.LEFT,
    word_wrap=True,
):
    return _add_textbox(
        slide,
        left_mm,
        top_mm,
        w_mm,
        h_mm,
        text,
        _fit_pt(text, base_pt, min_pt, w_mm - 0.8, h_mm - 0.3),
        bold=bold,
        color=color,
        align=align,
        word_wrap=word_wrap,
    )


def _label_and_value(
    slide, left_mm, top_mm, w_mm, h_mm, label, value, label_fs=None, value_fs=None
):
    label_fs = label_fs or tc.FS_LABEL
    value_fs = value_fs or tc.FS_VALUE
    label_h = h_mm * 0.32
    value_h = h_mm * 0.68
    _add_textbox(
        slide,
        left_mm + 0.35,
        top_mm,
        w_mm - 0.7,
        label_h,
        label,
        label_fs,
        color=tc.MID_GREY,
    )
    _fitted_textbox(
        slide,
        left_mm + 0.35,
        top_mm + label_h,
        w_mm - 0.7,
        value_h,
        value,
        value_fs.pt,
        min_pt=6.5,
        bold=True,
        color=tc.BLACK,
        align=PP_ALIGN.CENTER,
        word_wrap=True,
    )


def draw_titleblock(slide, fields: dict, logo_path: str = None):
    def fv(key, default=""):
        return str(fields.get(key, default) or default)

    company_name = fv("company_name", "PTG CONSULTING")
    company_address = fv(
        "company_address",
        "Level 3, 159 Coronation Drive (CNR Cribb St), Milton QLD 4064",
    )
    company_phone = fv("company_phone", "(07) 3444 6666")
    company_email = fv("company_email", "admin@ptgconsulting.com.au")
    company_website = fv("company_website", "www.ptgconsulting.com.au")

    _add_rect(
        slide,
        tc.MARGIN_L,
        tc.SLIDE_H_MM - tc.MARGIN_T,
        tc.MARGIN_R - tc.MARGIN_L,
        tc.MARGIN_T - tc.MARGIN_B,
        line_pt=0.25,
    )
    _add_rect(
        slide,
        tc.INNER_L,
        tc.SLIDE_H_MM - tc.INNER_T,
        tc.INNER_R - tc.INNER_L,
        tc.INNER_T - tc.INNER_B,
        line_pt=0.75,
    )

    tb_top_ppt = tc.SLIDE_H_MM - tc.TB_TOP_DXF
    tb_h = tc.TB_H_MM
    _add_rect(
        slide,
        tc.INNER_L,
        tb_top_ppt,
        tc.INNER_R - tc.INNER_L,
        tb_h,
        fill=tc.WHITE,
        line_pt=0.5,
    )

    meta_prepared_end = 718.86
    meta_approved_end = 748.86
    meta_status_end = 782.86
    meta_sheet_end = 809.86

    for x in [
        tc.COL_REV_END,
        tc.COL_COPY_END,
        tc.COL_ADDR_END,
        tc.COL_NORTH_END,
        tc.COL_PROJ_END,
        tc.COL_TITLE_END,
        meta_prepared_end,
        meta_approved_end,
        meta_status_end,
        meta_sheet_end,
    ]:
        _add_line(slide, x, tb_top_ppt, x, tb_top_ppt + tb_h, 0.4)

    row_top_ppt = tc.SLIDE_H_MM - tc.ROW_TOP_DXF
    row_mid_ppt = tc.SLIDE_H_MM - tc.ROW_MID_DXF
    _add_line(slide, tc.COL_TITLE_END, row_top_ppt, tc.INNER_R, row_top_ppt, 0.4)
    _add_line(slide, tc.COL_TITLE_END, row_mid_ppt, meta_status_end, row_mid_ppt, 0.4)

    for label, x in [
        ("REV", tc.INNER_L),
        ("DATE", tc.REV_COL_DATE),
        ("DESCRIPTION", tc.REV_COL_DESC),
        ("DRAWN", tc.REV_COL_DRAWN),
        ("CHK", tc.REV_COL_CHK),
    ]:
        y_ppt = tc.SLIDE_H_MM - (tc.REV_ROW_END + 6)
        _add_textbox(
            slide, x + 0.5, y_ppt, 20, 5, label, tc.FS_LABEL, color=tc.MID_GREY
        )
    for i in range(13):
        y = tc.REV_ROW_START + i * tc.REV_ROW_STEP
        y_ppt = tc.SLIDE_H_MM - y
        _add_line(slide, tc.INNER_L, y_ppt, tc.COL_REV_END, y_ppt, 0.25)
    for x in [tc.REV_COL_DATE, tc.REV_COL_DESC, tc.REV_COL_DRAWN, tc.REV_COL_CHK]:
        y_top = tc.SLIDE_H_MM - tc.REV_ROW_END
        y_bot = tc.SLIDE_H_MM - tc.REV_ROW_START
        _add_line(slide, x, y_top, x, y_bot, 0.25)

    cx_l = tc.COL_REV_END
    cx_r = tc.COL_COPY_END
    copyright_text = (
        f"This document remains the property of {company_name} and may not be copied "
        "or reproduced without permission.\nDo not scale this drawing.\nVerify all "
        "dimensions on site."
    )
    _add_textbox(
        slide,
        cx_l + 1,
        tb_top_ppt + 1,
        cx_r - cx_l - 2,
        28,
        copyright_text,
        tc.FS_SMALL,
        color=tc.DARK_GREY,
        word_wrap=True,
    )

    ax_l = tc.COL_COPY_END
    ax_r = tc.COL_ADDR_END
    ax_w = ax_r - ax_l
    _fitted_textbox(
        slide,
        ax_l + 0.8,
        tb_top_ppt + 1.2,
        ax_w - 1.6,
        8.2,
        f"a: {company_address}",
        5.6,
        min_pt=4.8,
        color=tc.DARK_GREY,
    )
    _fitted_textbox(
        slide,
        ax_l + 0.8,
        tb_top_ppt + 9.3,
        ax_w - 1.6,
        4.5,
        f"p: {company_phone}",
        5.6,
        min_pt=4.8,
        color=tc.DARK_GREY,
        word_wrap=False,
    )
    _fitted_textbox(
        slide,
        ax_l + 0.8,
        tb_top_ppt + 14.0,
        ax_w - 1.6,
        4.5,
        f"e: {company_email}",
        5.6,
        min_pt=4.8,
        color=tc.DARK_GREY,
        word_wrap=False,
    )
    _fitted_textbox(
        slide,
        ax_l + 0.8,
        tb_top_ppt + 18.7,
        ax_w - 1.6,
        4.5,
        f"w: {company_website}",
        5.6,
        min_pt=4.8,
        color=tc.DARK_GREY,
        word_wrap=False,
    )

    logo_l = tc.COL_ADDR_END
    logo_r = tc.COL_NORTH_END
    logo_w = logo_r - logo_l
    if logo_path and os.path.isfile(logo_path):
        try:
            slide.shapes.add_picture(
                logo_path,
                _mm(logo_l + 2),
                _mm(tb_top_ppt + 2),
                _mm(logo_w - 4),
                _mm(tb_h - 4),
            )
        except Exception:
            _fitted_textbox(
                slide,
                logo_l + 1.2,
                tb_top_ppt + tb_h / 2 - 5.5,
                logo_w - 2.4,
                11,
                company_name,
                12.0,
                min_pt=9,
                bold=True,
                color=tc.PTG_BLUE,
                align=PP_ALIGN.CENTER,
            )
    else:
        _fitted_textbox(
            slide,
            logo_l + 1.2,
            tb_top_ppt + tb_h / 2 - 5.5,
            logo_w - 2.4,
            11,
            company_name,
            12.5,
            min_pt=9.5,
            bold=True,
            color=tc.PTG_BLUE,
            align=PP_ALIGN.CENTER,
        )

    cl_l = tc.COL_NORTH_END
    cl_r = tc.COL_PROJ_END
    cl_w = cl_r - cl_l
    _add_textbox(
        slide,
        cl_l + 0.8,
        tb_top_ppt + 0.5,
        cl_w - 1.6,
        4.2,
        "Client",
        tc.FS_LABEL,
        color=tc.MID_GREY,
    )
    _fitted_textbox(
        slide,
        cl_l + 0.8,
        tb_top_ppt + 4.4,
        cl_w - 1.6,
        14,
        fv("client_name"),
        10.5,
        min_pt=7.2,
        bold=True,
    )

    pt_l = tc.COL_PROJ_END
    pt_r = tc.COL_TITLE_END
    pt_w = pt_r - pt_l
    _add_textbox(
        slide,
        pt_l + 0.8,
        tb_top_ppt + 0.5,
        pt_w - 1.6,
        4,
        "Project",
        tc.FS_LABEL,
        color=tc.MID_GREY,
    )
    _fitted_textbox(
        slide,
        pt_l + 0.8,
        tb_top_ppt + 3.8,
        pt_w - 1.6,
        14.6,
        fv("project_title"),
        tc.FS_PROJECT.pt,
        min_pt=8.5,
        bold=True,
        color=tc.PTG_BLUE,
    )
    _add_textbox(
        slide,
        pt_l + 0.8,
        tb_top_ppt + 18.2,
        pt_w - 1.6,
        4,
        "Address",
        tc.FS_LABEL,
        color=tc.MID_GREY,
    )
    _fitted_textbox(
        slide,
        pt_l + 0.8,
        tb_top_ppt + 21.9,
        pt_w - 1.6,
        7.2,
        fv("project_address"),
        8.4,
        min_pt=6.2,
    )
    _add_textbox(
        slide,
        pt_l + 0.8,
        tb_top_ppt + 29.2,
        18,
        3.6,
        "Project No.",
        tc.FS_LABEL,
        color=tc.MID_GREY,
    )
    _fitted_textbox(
        slide,
        pt_l + 18.2,
        tb_top_ppt + 28.9,
        pt_w - 19.0,
        4.2,
        fv("project_number"),
        8.2,
        min_pt=6.5,
        bold=True,
        word_wrap=False,
    )

    dt_l = tc.COL_TITLE_END
    dt_r = meta_prepared_end
    dt_w = dt_r - dt_l
    _add_textbox(
        slide,
        dt_l + 0.8,
        tb_top_ppt + 0.5,
        dt_w - 1.6,
        4,
        "Drawing Title",
        tc.FS_LABEL,
        color=tc.MID_GREY,
    )
    _fitted_textbox(
        slide,
        dt_l + 0.8,
        tb_top_ppt + 3.8,
        dt_w - 1.6,
        14.6,
        fv("drawing_title_1"),
        tc.FS_TITLE.pt,
        min_pt=8.5,
        bold=True,
    )
    if fv("drawing_title_2"):
        _fitted_textbox(
            slide,
            dt_l + 0.8,
            tb_top_ppt + 18.2,
            dt_w - 1.6,
            8.4,
            fv("drawing_title_2"),
            9.5,
            min_pt=7.0,
            bold=True,
        )
    if fv("drawing_title_3"):
        _fitted_textbox(
            slide,
            dt_l + 0.8,
            tb_top_ppt + 26.1,
            dt_w - 1.6,
            6.4,
            fv("drawing_title_3"),
            8.2,
            min_pt=6.4,
        )

    sub_top_h = row_top_ppt - tb_top_ppt
    _label_and_value(
        slide,
        tc.COL_TITLE_END,
        tb_top_ppt,
        meta_prepared_end - tc.COL_TITLE_END,
        sub_top_h,
        "Drawn",
        fv("drawn_by"),
    )
    _label_and_value(
        slide,
        meta_prepared_end,
        tb_top_ppt,
        meta_approved_end - meta_prepared_end,
        sub_top_h,
        "Approved",
        fv("approved_by"),
    )
    _label_and_value(
        slide,
        meta_approved_end,
        tb_top_ppt,
        meta_status_end - meta_approved_end,
        sub_top_h,
        "Scale @ A1",
        fv("scale"),
    )
    _label_and_value(
        slide,
        meta_status_end,
        tb_top_ppt,
        meta_sheet_end - meta_status_end,
        sub_top_h,
        "Figure No.",
        fv("sheet_prefix", "A") + fv("sheet_number", "001"),
    )
    _label_and_value(
        slide,
        meta_sheet_end,
        tb_top_ppt,
        tc.INNER_R - meta_sheet_end,
        sub_top_h,
        "Revision",
        fv("revision", "-"),
    )

    sub_mid_h = row_mid_ppt - row_top_ppt
    _label_and_value(
        slide,
        tc.COL_TITLE_END,
        row_top_ppt,
        meta_prepared_end - tc.COL_TITLE_END,
        sub_mid_h,
        "Designed",
        fv("designed_by"),
    )
    _label_and_value(
        slide,
        meta_prepared_end,
        row_top_ppt,
        meta_approved_end - meta_prepared_end,
        sub_mid_h,
        "Date",
        fv("date"),
        value_fs=Pt(8.3),
    )
    _label_and_value(
        slide,
        meta_approved_end,
        row_top_ppt,
        meta_status_end - meta_approved_end,
        sub_mid_h,
        "Drawing Status",
        fv("drawing_status", "FOR INFORMATION"),
        value_fs=Pt(8.0),
    )

    status = fv("drawing_status", "FOR INFORMATION").upper()
    if "NOT FOR CONSTRUCTION" in status or "DRAFT" in status:
        _add_textbox(
            slide,
            tc.CONTENT_L + tc.CONTENT_W * 0.3,
            tc.CONTENT_T + tc.CONTENT_H * 0.35,
            tc.CONTENT_W * 0.4,
            30,
            "DRAFT",
            Pt(72),
            bold=True,
            color=RGBColor(0xCC, 0xCC, 0xCC),
            align=PP_ALIGN.CENTER,
        )
