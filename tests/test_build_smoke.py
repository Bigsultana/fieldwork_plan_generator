import sys
import tempfile
import unittest
from pathlib import Path

from PIL import Image
from pptx import Presentation

ROOT = Path(__file__).resolve().parents[1]
SOURCE = ROOT / "source"
if str(SOURCE) not in sys.path:
    sys.path.insert(0, str(SOURCE))

from utils.builder import build_appendix
from utils.config_loader import validate_sheet_config


class BuildAppendixSmokeTest(unittest.TestCase):
    def test_build_appendix_generates_generic_presentation(self):
        with tempfile.TemporaryDirectory() as temporary_directory:
            directory = Path(temporary_directory)
            image_path = directory / "sample_plan.png"
            output_path = directory / "appendix.pptx"
            Image.new("RGB", (1600, 900), "white").save(image_path)

            result = build_appendix(
                exports_dir=str(directory),
                output_path=str(output_path),
                project_config={
                    "project_title": "Smoke Test Project",
                    "project_address": "123 Test Street",
                    "project_number": "TEST-001",
                    "client_name": "Example Client",
                    "drawn_by": "AB",
                    "designed_by": "CD",
                    "approved_by": "EF",
                    "date": "30.07.26",
                    "drawing_status": "FOR INFORMATION",
                    "sheet_prefix": "A",
                    "company_name": "Example Company",
                    "company_address": "Example Address",
                    "company_phone": "",
                    "company_email": "",
                    "company_website": "",
                },
                sheets=[
                    {
                        "sheet_number": "001",
                        "drawing_title_1": "Sample Plan",
                        "drawing_title_2": "",
                        "drawing_title_3": "",
                        "filename_pattern": "sample_plan",
                        "scale": "NTS",
                    }
                ],
            )

            self.assertEqual(result["slides_built"], 1)
            self.assertEqual(result["missing"], [])
            self.assertTrue(output_path.is_file())
            presentation = Presentation(str(output_path))
            self.assertEqual(len(presentation.slides), 1)
            text = "\n".join(
                shape.text
                for shape in presentation.slides[0].shapes
                if hasattr(shape, "text_frame")
            )
            self.assertIn("Example Company", text)

    def test_duplicate_sheet_numbers_are_rejected(self):
        errors = validate_sheet_config(
            [
                {
                    "sheet_number": "001",
                    "drawing_title_1": "One",
                    "filename_pattern": "one",
                },
                {
                    "sheet_number": "001",
                    "drawing_title_1": "Two",
                    "filename_pattern": "two",
                },
            ]
        )
        self.assertTrue(
            any("duplicate sheet_number" in error for error in errors)
        )


if __name__ == "__main__":
    unittest.main()
