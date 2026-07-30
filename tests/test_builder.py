import sys
import tempfile
import unittest
from pathlib import Path

from PIL import Image
from pptx import Presentation

ROOT = Path(__file__).resolve().parents[1]
SOURCE = ROOT / "source"
if str(SOURCE) not in sys.path:
    sys.path.insert(0, str(SOURCE))

from utils.builder import build_fieldwork_plan
from utils.config_loader import clean_project_config, validate_sheet_config


class FieldworkPlanBuilderTest(unittest.TestCase):
    def test_builds_a_generic_fieldwork_plan(self):
        with tempfile.TemporaryDirectory() as temporary_directory:
            directory = Path(temporary_directory)
            image_path = directory / "location-plan.png"
            output_path = directory / "fieldwork-plan.pptx"
            Image.new("RGB", (1600, 900), "white").save(image_path)

            result = build_fieldwork_plan(
                exports_dir=str(directory),
                output_path=str(output_path),
                project_config=clean_project_config(
                    {
                        "project_title": "Test Fieldwork Plan",
                        "project_number": "TEST-001",
                        "client_name": "Example Client",
                        "company_name": "Example Company",
                    }
                ),
                sheets=[
                    {
                        "sheet_number": "001",
                        "drawing_title_1": "Test Location Plan",
                        "drawing_title_2": "Fieldwork access and proposed locations",
                        "scale": "NTS",
                        "revision": "A",
                        "source_path": str(image_path),
                    }
                ],
            )

            self.assertEqual(result["slides_built"], 1)
            self.assertEqual(result["missing"], [])
            presentation = Presentation(str(output_path))
            self.assertEqual(len(presentation.slides), 1)
            text = "\n".join(
                shape.text
                for shape in presentation.slides[0].shapes
                if hasattr(shape, "text_frame")
            )
            self.assertIn("Example Company", text)
            self.assertIn("Test Location Plan", text)

    def test_duplicate_sheet_numbers_are_rejected(self):
        errors = validate_sheet_config(
            [
                {
                    "sheet_number": "001",
                    "drawing_title_1": "One",
                    "source_path": "one.png",
                },
                {
                    "sheet_number": "001",
                    "drawing_title_1": "Two",
                    "source_path": "two.png",
                },
            ]
        )
        self.assertTrue(any("duplicate sheet number" in error for error in errors))


if __name__ == "__main__":
    unittest.main()
